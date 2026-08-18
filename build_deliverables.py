"""Build DA1_Report.docx and DA1_Presentation.pptx from the diagram PNGs.

Run: python build_deliverables.py     (run diagrams/make_diagrams.py first)
"""
import os
import re
import shutil
import zipfile
from datetime import datetime
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = os.path.dirname(os.path.abspath(__file__))
DIAG = os.path.join(ROOT, "diagrams")

BLACK = RGBColor(0x00, 0x00, 0x00)
DGREY = RGBColor(0x33, 0x33, 0x33)
PBLACK = PRGB(0x00, 0x00, 0x00)
PGREY = PRGB(0x44, 0x44, 0x44)
PWHITE = PRGB(0xFF, 0xFF, 0xFF)
PHEAD = PRGB(0xD9, 0xD9, 0xD9)
PZEBRA = PRGB(0xF4, 0xF4, 0xF4)

TITLE = "Warehouse and Supply Chain Management System"
TEAM = [("Safwaan Mohamed S", "25BCE1188"),
        ("Avi Dhandhania", "25BCE1207"),
        ("Priyal Bhalla", "25BCE1261")]

# document properties, so the files carry the team's own authorship
DOC_AUTHOR = "Avi Dhandhania"
CREATED = datetime(2026, 7, 21, 19, 42)
MODIFIED = datetime(2026, 7, 28, 22, 15)

ENTITIES = [
    ("WAREHOUSE", "WarehouseID", "WName, City, Capacity"),
    ("PRODUCT", "ProductID", "PName, Category, UnitPrice"),
    ("SUPPLIER", "SupplierID", "SName, City, Phone (many per supplier)"),
    ("CUSTOMER", "CustomerID", "CName, Address (Street + City + Pincode), Phone"),
    ("ORDERS", "OrderID", "OrderDate, Status, TotalAmt (calculated)"),
    ("SHIPMENT", "ShipmentID", "DispatchDate, DeliveryDate, Status"),
    ("GOODS_RETURN", "ReturnID", "ReturnDate, Reason, RefundAmt, Status"),
    ("INSPECTION", "InspID", "InspDate, Result"),
    ("EMPLOYEE", "EmpID", "EName, Role, Salary"),
]

RELATIONSHIPS = [
    ("SUPPLIES", "SUPPLIER - PRODUCT", "M:N", "A supplier supplies many products, "
     "and a product comes from many suppliers. Carries LeadTimeDays."),
    ("STORED_IN", "PRODUCT - WAREHOUSE", "M:N", "A product is stocked in many "
     "warehouses. Carries Quantity."),
    ("CONTAINS", "PRODUCT - ORDERS", "M:N", "An order contains many products. "
     "Carries Qty."),
    ("WORKS_AT", "WAREHOUSE - EMPLOYEE", "1:N", "Every employee is posted to "
     "exactly one warehouse."),
    ("PLACES", "CUSTOMER - ORDERS", "1:N", "Every order belongs to exactly one "
     "customer."),
    ("SHIPPED_BY", "ORDERS - SHIPMENT", "1:N", "One order can go out in more than "
     "one shipment."),
    ("DISPATCHED_FROM", "WAREHOUSE - SHIPMENT", "1:N", "Every shipment leaves from "
     "one warehouse."),
    ("MANAGES", "EMPLOYEE - WAREHOUSE", "1:1", "A warehouse has exactly one "
     "manager, and an employee manages at most one warehouse. The only 1:1 "
     "relationship in the design."),
    ("RETURNED_FOR", "ORDERS - GOODS_RETURN", "1:N", "A return is always against "
     "one order, and an order can be returned in parts."),
    ("RETURNED_TO", "WAREHOUSE - GOODS_RETURN", "1:N", "Returned goods come back "
     "into one warehouse."),
    ("CHECKS_SUPPLIER", "INSPECTION - SUPPLIER", "N:1", "Names the supplier the "
     "inspection is about. Replaced by the aggregation in the EER model."),
    ("CHECKS_PRODUCT", "INSPECTION - PRODUCT", "N:1", "Names the product the "
     "inspection is about. Replaced by the aggregation in the EER model."),
]

EER_ADDED = [
    ("Weak entity", "ORDER_ITEM", "An order line means nothing without its order, "
     "and ItemNo (1, 2, 3 and so on) is only unique inside one order, so its key is "
     "OrderID + ItemNo. This replaces the M:N CONTAINS relationship."),
    ("Specialization\n(total, disjoint)", "EMPLOYEE to\nMANAGER, DRIVER",
     "Every employee is exactly one of the two, never both and never neither. "
     "A manager has Level and Bonus, a driver has LicenseNo and Expiry."),
    ("Specialization\n(partial, overlapping)",
     "PRODUCT to\nPERISHABLE_PRODUCT,\nHAZARDOUS_PRODUCT",
     "Only some products are perishable and only some are hazardous, so the "
     "specialization is partial: a plain product joins neither subclass. It is also "
     "overlapping, because a product can be both, a vaccine being the obvious case. "
     "That is what the o circle here says and the d circles elsewhere deny."),
    ("Multivalued attribute", "SUPPLIER.Phone",
     "One supplier can have several phone numbers."),
    ("Composite attribute", "CUSTOMER.Address",
     "Address is made of Street, City and Pincode, which we sometimes need "
     "separately."),
    ("Derived attribute", "ORDERS.TotalAmt",
     "It can always be worked out by adding up the order lines, so storing it "
     "would only create a second copy that can go stale."),
    ("Recursive relationship", "EMPLOYEE\nSUPERVISES EMPLOYEE",
     "An employee reports to another employee, so the relationship joins one "
     "entity to itself."),
    ("Generalization\n(total, disjoint)",
     "SHIPMENT + GOODS_RETURN\nto GOODS_MOVEMENT",
     "Both entities were already in the ER model, and both turned out to be one "
     "order's goods moving on a date with a status. Those shared attributes, and "
     "the links to ORDERS and WAREHOUSE, are lifted into a common superclass. "
     "Built upwards from what existed, which is what makes it generalization "
     "rather than specialization. ShipmentID and ReturnID collapse into one key, "
     "MovementID."),
    ("Aggregation", "SUPPLIER SUPPLIES PRODUCT,\nthen INSPECTED_IN",
     "An inspection is of one supplier's supply of one product, not of the "
     "supplier and not of the product. Boxing the whole SUPPLIES relationship and "
     "treating it as one object lets INSPECTION point at the pair, replacing the "
     "two separate links the ER model needed."),
    ("1:1 relationship\nmade precise", "WAREHOUSE MANAGES\nto MANAGER",
     "The ER model can only say a warehouse is managed by an employee. Once "
     "EMPLOYEE is specialized, the same 1:1 relationship joins WAREHOUSE to "
     "MANAGER, which is what the business rule actually means."),
]

MAPPING_RULES = [
    "Each entity becomes one table, and its key attribute becomes the primary key.",
    "Each 1:N relationship becomes a foreign key stored on the 'many' side.",
    "Each M:N relationship becomes a new table whose primary key is the two keys "
    "put together.",
    "Each multivalued attribute becomes its own small table.",
    "Each composite attribute becomes one column per part.",
    "A derived attribute is not stored at all, it is calculated when needed.",
    "A weak entity becomes a table keyed by its owner's key plus its own partial key.",
    "A specialization becomes one table per subclass, each sharing the "
    "superclass's key.",
    "A generalization is mapped the same way: one table for the superclass and "
    "one for each subclass, each subclass keyed by the superclass's key.",
    "Each 1:1 relationship becomes a foreign key on the side where taking part is "
    "compulsory, with a UNIQUE constraint so that one row cannot be claimed twice.",
    "An aggregation becomes a foreign key pointing at the primary key of the table "
    "the aggregated relationship produced.",
]

SCHEMA = [
    ("CITY", "City", "none"),
    ("SUPPLIER", "SupplierID", "City"),
    ("SUPPLIER_PHONE", "SupplierID + Phone", "SupplierID"),
    ("PRODUCT", "ProductID", "none"),
    ("PERISHABLE_PRODUCT", "ProductID", "ProductID"),
    ("HAZARDOUS_PRODUCT", "ProductID", "ProductID"),
    ("SUPPLIES", "SupplierID + ProductID", "SupplierID, ProductID"),
    ("INSPECTION", "InspID", "SupplierID + ProductID (to SUPPLIES)"),
    ("WAREHOUSE", "WarehouseID", "City, ManagerID (UNIQUE)"),
    ("EMPLOYEE", "EmpID", "WarehouseID, SupervisorID"),
    ("MANAGER", "EmpID", "EmpID"),
    ("DRIVER", "EmpID", "EmpID"),
    ("DRIVER_CITY", "DriverID", "DriverID, City"),
    ("STOCK", "WarehouseID + ProductID", "WarehouseID, ProductID"),
    ("CUSTOMER", "CustID", "City"),
    ("ORDERS", "OrderID", "CustID"),
    ("ORDER_ITEM", "OrderID + ItemNo", "OrderID, ProductID"),
    ("ORDER_DRIVER", "OrderID + DriverID", "OrderID, DriverID"),
    ("GOODS_MOVEMENT", "MovementID", "OrderID"),
    ("SHIPMENT", "MovementID", "MovementID, WarehouseID"),
    ("GOODS_RETURN", "MovementID", "MovementID"),
]

# ---- dependencies read off the sample rows in Figure 3
# (ref, dependency, evidence for the report, short version for the slide)
FD_HOLD = [
    ("F1", "OrderID -> OrderDate, Status, CustID",
     "Rows 1 and 2 are both O-101 and agree on all three columns, and rows 6 and 7 "
     "do the same for O-105.",
     "O-101 in rows 1 and 2, O-105 in rows 6 and 7"),
    ("F2", "CustID -> CName, Street, City",
     "C-01 appears in rows 1, 2 and 4 with the same name, street and city, and C-03 "
     "appears in rows 3, 6 and 7.",
     "C-01 in rows 1, 2, 4; C-03 in rows 3, 6, 7"),
    ("F3", "City -> State",
     "Chennai always carries Tamil Nadu (rows 1, 2, 4, 5) and Mumbai always carries "
     "Maharashtra (rows 3, 6, 7).",
     "Chennai is always Tamil Nadu, Mumbai always Maharashtra"),
    ("F4", "ProductID -> PName, Category, Price",
     "P-100 gives 480 in rows 1 and 4, P-104 gives 132 in rows 2 and 6, P-210 gives "
     "96 in rows 3 and 7.",
     "P-100 is 480 twice, P-104 is 132 twice, P-210 is 96 twice"),
    ("F5", "{OrderID, ProductID} -> Qty",
     "The pair is unique across the seven rows, and neither column alone fixes Qty "
     "(see N1 and N2).",
     "the pair is unique, and neither half works alone"),
    ("F6", "DriverID -> City",
     "In DELIVERY_DUTY, D-05 is Chennai in rows 1, 4 and 5, and D-09 is Mumbai in "
     "rows 2, 3 and 6.",
     "D-05 is always Chennai, D-09 always Mumbai"),
    ("F7", "{OrderID, City} -> DriverID",
     "No two rows of DELIVERY_DUTY share both an OrderID and a City, so the pair "
     "fixes the driver.",
     "no two rows share an OrderID and a City"),
    ("F8", "Phone -> SupplierID",
     "Every number in SUPPLIER_PHONE appears against one supplier only, so Phone is "
     "a second candidate key of that relation.",
     "each number belongs to one supplier"),
]

FD_REJECT = [
    ("N1", "ProductID -> Qty",
     "Rows 1 and 4: P-100 is ordered 20 times on O-101 and 15 times on O-103."),
    ("N2", "OrderID -> ProductID",
     "Rows 1 and 2: O-101 holds both P-100 and P-104."),
    ("N3", "City -> CustID",
     "Rows 1 and 5: Chennai holds C-01 and also C-05."),
    ("N4", "Category -> ProductID",
     "Rows 1 and 5: Staples covers P-100 and also P-108."),
    ("N5", "CName -> OrderID",
     "Rows 1 and 4: Sharma Retail placed O-101 and also O-103."),
    ("N6", "City -> DriverID",
     "DELIVERY_DUTY rows 1 and 7: Chennai is covered by D-05 and by D-11."),
    ("N7", "OrderID -> DriverID",
     "DELIVERY_DUTY rows 1 and 2: O-101 is served by D-05 and by D-09."),
    ("N8", "SupplierID -> Phone",
     "SUPPLIER_PHONE rows 1 and 2: S-01 has two numbers."),
]

NORM_STEPS = [
    ("norm_unf.png", "8.1  Where we start, Unnormalized Form (UNF)",
     "Everything sits in one wide table, the way a clerk would keep it in a "
     "spreadsheet. One cell holds several phone numbers, another holds a whole group "
     "of product lines, and a third holds the drivers the order went out with. "
     "Because the customer's details are copied onto every order row, the same fact "
     "is written down many times, and that is what causes the insert, delete and "
     "update problems described in section 1.1.",
     "Repeating groups: SupplierPhones, Items and Delivery. Because those cells are "
     "not atomic there is no key yet, so no dependency can even be stated.",
     ["One wide table, the way a clerk keeps it in a spreadsheet.",
      "Three cells hold lists: SupplierPhones, Items and Delivery.",
      "Customer details are copied onto every order row.",
      "No atomic key, so no dependency can even be written down yet."]),
    ("norm_1nf.png", "8.2  First Normal Form (1NF)",
     "1NF asks for two things: every cell must hold exactly one value, and a "
     "repeating group must move into its own table. The product lines become "
     "ORDER_ITEM, the phone numbers become SUPPLIER_PHONE, and the delivery pairs "
     "become DELIVERY_DUTY. ORDER_ITEM needs OrderID and ProductID together as its "
     "key, because neither column alone identifies a row (N1 and N2). DELIVERY_DUTY "
     "is keyed by OrderID and City, from F7.",
     "Still wrong: PName and Category repeat on every order line that mentions the "
     "same product, which is F4 sitting inside a table whose key is "
     "{OrderID, ProductID}.",
     ["Rule: one value per cell, and repeating groups move out.",
      "Items become ORDER_ITEM, phones become SUPPLIER_PHONE, deliveries become "
      "DELIVERY_DUTY.",
      "ORDER_ITEM is keyed by OrderID + ProductID, because N1 and N2 rule out "
      "either column alone.",
      "DELIVERY_DUTY is keyed by OrderID + City, from F7.",
      "Still wrong: F4 means PName and Category repeat on every line for the same "
      "product."]),
    ("norm_2nf.png", "8.3  Second Normal Form (2NF)",
     "2NF removes partial dependency, which means a column must depend on the whole "
     "key and not on part of it. ORDER_ITEM is keyed by {OrderID, ProductID}, but F4 "
     "says PName, Category and Price are fixed by ProductID alone. They describe the "
     "product and not the order, so they move to a new PRODUCT table where Price "
     "becomes UnitPrice. Only Qty depends on the whole key (F5), so only Qty stays. "
     "The other three tables have no partial dependency and are left alone.",
     "Still wrong: in ORDER_MASTER the customer's name is reached through CustID (F2) "
     "and the state is reached through the city (F3). Neither route starts at the "
     "key.",
     ["Rule: no column may depend on only part of the key.",
      "ORDER_ITEM is keyed by {OrderID, ProductID}, but F4 fixes PName, Category "
      "and Price from ProductID alone.",
      "Those three move to PRODUCT, where Price becomes UnitPrice.",
      "Only Qty depends on the whole key (F5), so only Qty stays behind.",
      "Still wrong: F2 and F3 in ORDER_MASTER both start from a non-key column."]),
    ("norm_3nf.png", "8.4  Third Normal Form (3NF)",
     "3NF removes transitive dependency, which means one ordinary column must not "
     "decide another. ORDER_MASTER is keyed by OrderID, yet F2 and F3 both start "
     "from a non-key column, so each determinant is pulled out into a table of its "
     "own: CUSTOMER keyed by CustID and CITY keyed by City. What is left of "
     "ORDER_MASTER becomes ORDERS, holding only what OrderID decides directly, plus "
     "CustID as a foreign key.",
     "One table still hides a problem. DELIVERY_DUTY satisfies 3NF, because every "
     "column in it belongs to some candidate key, and 3NF excuses a non-key "
     "determinant on exactly that ground.",
     ["Rule: no ordinary column may decide another ordinary column.",
      "F2 (CustID -> CName, Street, City) becomes CUSTOMER.",
      "F3 (City -> State) becomes CITY.",
      "What is left of ORDER_MASTER becomes ORDERS, keeping CustID as a foreign key.",
      "Still wrong: DELIVERY_DUTY passes 3NF only because every column of it "
      "belongs to some key."]),
    ("norm_bcnf.png", "8.5  Boyce-Codd Normal Form (BCNF)",
     "BCNF applies one strict rule: the left-hand side of every dependency must be a "
     "key. DELIVERY_DUTY has two candidate keys, {OrderID, City} and "
     "{OrderID, DriverID}, so F7 is fine. F6 is not, because DriverID decides City "
     "while not being a key on its own. That is why 'D-09 covers Mumbai' is stored on "
     "three separate rows of the sample data and why two of those rows can be made to "
     "disagree. Splitting on F6 gives DRIVER_CITY, keyed by DriverID, and "
     "ORDER_DRIVER. Joining the two back on DriverID returns the original seven rows "
     "exactly, so the split loses nothing.",
     "Every table is now in BCNF, which also means every table is in 1NF, 2NF and "
     "3NF.",
     ["Rule: the left side of every dependency must be a key.",
      "DELIVERY_DUTY has two candidate keys, {OrderID, City} and "
      "{OrderID, DriverID}, so F7 is fine.",
      "F6 (DriverID -> City) is not, because DriverID is not a key on its own.",
      "'D-09 covers Mumbai' therefore sits on three rows, and two of them can be "
      "made to disagree.",
      "Splitting on F6 gives DRIVER_CITY and ORDER_DRIVER, and the join on DriverID "
      "returns the original rows.",
      "Every table is now in BCNF."]),
]


def img(name):
    return os.path.join(DIAG, name)


def stamp_office_metadata(path, app, fmt=None):
    """Fix up the extended properties left behind by the template.

    python-docx and python-pptx ship a blank file that was made in Mac Office
    2011, so every document they write claims to be that, and carries that
    template's thumbnail of an empty page. Both are rewritten here. The same
    template also pins Word's compatibility mode to 2010, which is what puts
    'Compatibility Mode' in the title bar, so that is lifted too.
    """
    tmp = path + ".tmp"
    with zipfile.ZipFile(path) as zin, \
            zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            if item.filename == "docProps/thumbnail.jpeg":
                continue
            data = zin.read(item.filename)
            if item.filename == "docProps/app.xml":
                text = data.decode("utf-8")
                text = re.sub(r"<Application>[^<]*</Application>",
                              f"<Application>{app}</Application>", text)
                text = re.sub(r"<AppVersion>[^<]*</AppVersion>",
                              "<AppVersion>16.0000</AppVersion>", text)
                if fmt:
                    text = re.sub(r"<PresentationFormat>[^<]*</PresentationFormat>",
                                  f"<PresentationFormat>{fmt}</PresentationFormat>",
                                  text)
                data = text.encode("utf-8")
            elif item.filename == "_rels/.rels":
                data = re.sub(rb"<Relationship[^>]*thumbnail[^>]*/>", b"", data)
            elif item.filename == "word/settings.xml":
                data = data.replace(b'w:name="compatibilityMode" '
                                    b'w:uri="http://schemas.microsoft.com/office/word"'
                                    b' w:val="14"',
                                    b'w:name="compatibilityMode" '
                                    b'w:uri="http://schemas.microsoft.com/office/word"'
                                    b' w:val="15"')
            zout.writestr(item, data)
    shutil.move(tmp, path)
    return path


# =====================================================================  DOCX
def build_docx():
    d = Document()
    st = d.styles["Normal"]
    st.font.name = "Calibri"
    st.font.size = Pt(11)
    for s in d.sections:
        s.top_margin = s.bottom_margin = Inches(0.8)
        s.left_margin = s.right_margin = Inches(0.7)

    cp = d.core_properties
    cp.title = TITLE
    cp.subject = "Database Management Systems, Digital Assignment 1"
    cp.author = DOC_AUTHOR
    cp.last_modified_by = DOC_AUTHOR
    cp.category = ""
    cp.comments = ""
    cp.keywords = ""
    cp.created = CREATED
    cp.modified = MODIFIED
    cp.revision = 9

    def h(text, level):
        p = d.add_heading(text, level)
        for r in p.runs:
            r.font.color.rgb = BLACK
        return p

    def para(text="", bold=False, italic=False, size=11, align=None):
        p = d.add_paragraph()
        if align is not None:
            p.alignment = align
        r = p.add_run(text)
        r.bold = bold
        r.italic = italic
        r.font.size = Pt(size)
        return p

    def bullets(items):
        for it in items:
            p = d.add_paragraph(style="List Bullet")
            p.add_run(it).font.size = Pt(10.5)

    def table(headers, rows, widths=None, fs=9.5):
        t = d.add_table(rows=1, cols=len(headers))
        t.style = "Table Grid"
        for i, hd in enumerate(headers):
            c = t.rows[0].cells[i]
            c.text = ""
            r = c.paragraphs[0].add_run(hd)
            r.bold = True
            r.font.size = Pt(fs)
        for row in rows:
            cells = t.add_row().cells
            for i, val in enumerate(row):
                cells[i].text = ""
                cells[i].paragraphs[0].add_run(str(val)).font.size = Pt(fs)
        if widths:
            for row in t.rows:
                for i, w in enumerate(widths):
                    row.cells[i].width = Inches(w)
        return t

    def figure(name, caption, width=7.0):
        d.add_picture(img(name), width=Inches(width))
        d.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
        return para(caption, italic=True, size=9, align=WD_ALIGN_PARAGRAPH.CENTER)

    # ---------------- title page
    para()
    para(TITLE, bold=True, size=26, align=WD_ALIGN_PARAGRAPH.CENTER)
    para("Digital Assignment 1 (DA1)", size=15, italic=True,
         align=WD_ALIGN_PARAGRAPH.CENTER)
    para("ER / EER Model and Normalization up to BCNF", size=13, bold=True,
         align=WD_ALIGN_PARAGRAPH.CENTER)
    para()
    para("Database Management Systems", size=12, align=WD_ALIGN_PARAGRAPH.CENTER)
    para("Vellore Institute of Technology, Chennai", size=12,
         align=WD_ALIGN_PARAGRAPH.CENTER)
    para()
    para("Submitted by", size=12, italic=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    for name, reg in TEAM:
        para(f"{name} ({reg})", size=13, bold=True,
             align=WD_ALIGN_PARAGRAPH.CENTER)
    para()
    para("Submission deadline: 31 July 2026", size=11, italic=True,
         align=WD_ALIGN_PARAGRAPH.CENTER)
    d.add_page_break()

    # ---------------- 1
    h("1. Problem Statement, Objectives and Scope", 1)
    h("1.1 Problem Statement", 2)
    para("A warehouse business has to know three things at all times: what stock it "
         "holds and where, who supplies each product, and where every customer order "
         "has reached. In most small and mid-sized operations these three answers "
         "live in separate places, namely a purchase register, a stock sheet on the "
         "warehouse floor, and a delivery log kept by the transport team. Nothing "
         "keeps the three in agreement.")
    para("The result is that the same fact ends up written down in several places. A "
         "customer's address appears on every order they have ever placed, and a "
         "product's name appears on every order line that includes it. When one copy "
         "is corrected and another is missed, the data starts contradicting itself, "
         "and the following problems follow directly:")
    bullets([
        "A new supplier or customer cannot be recorded until they appear on an order.",
        "Deleting the last order for a customer erases that customer completely.",
        "Correcting a product name or an address means editing many rows, and any "
        "row that is missed leaves the database in two minds.",
        "Nobody can say reliably how much of a product is on hand, because the "
        "figure is kept in more than one place.",
    ])
    para("These are known as insert, delete and update anomalies. They are not "
         "caused by missing data, they are caused by the same data being stored more "
         "than once. Normalization is the process that removes that duplication, and "
         "this assignment applies it step by step.")

    h("1.2 Objectives", 2)
    bullets([
        "O1. Identify the entities, attributes and relationships the system needs.",
        "O2. Draw an ER model for the design using Chen notation.",
        "O3. Extend that same model into an EER model, adding specialization, "
        "generalization, aggregation, a weak entity, and multivalued, composite and "
        "derived attributes.",
        "O4. Convert the EER model into relational tables using the standard "
        "mapping rules.",
        "O5. Populate the tables with sample rows and work out every functional "
        "dependency that holds, along with the candidate keys they imply.",
        "O6. Normalize the tables step by step from UNF through 1NF, 2NF and 3NF up "
        "to BCNF, showing the tables at each stage.",
        "O7. Confirm that no table in the final design can hold a contradiction.",
    ])

    h("1.3 Scope", 2)
    para("In scope: ", bold=True)
    d.paragraphs[-1].add_run(
        "nine entities covering warehouses, products, suppliers, customers, "
        "orders, shipments, returns, inspections and employees, the twelve "
        "relationships between them, the conversion of the design into tables, the "
        "dependency analysis, and normalization up to BCNF.")
    para("Out of scope for DA1: ", bold=True)
    d.paragraphs[-1].add_run(
        "writing the SQL and PL/SQL, which is DA2, and rebuilding the project on a "
        "modern database technology, which is DA3. Billing, tax, user logins and "
        "demand forecasting are left out of the design altogether so that the model "
        "stays small enough to explain in full.")

    # ---------------- 2
    h("2. Entities and Attributes", 1)
    para("The design uses nine entities. The key attribute of each one is the "
         "column that identifies a single row, and it is shown underlined in every "
         "diagram in this report.")
    table(["Entity", "Key attribute", "Other attributes"],
          ENTITIES, widths=[1.4, 1.3, 3.9])

    # ---------------- 3
    h("3. Relationships", 1)
    para("Twelve relationships connect the entities. '1:1' means one row on each "
         "side. '1:N' means one row on the left can be linked to many rows on the "
         "right. 'M:N' means many on both sides.")
    table(["Relationship", "Between", "Type", "Meaning"],
          RELATIONSHIPS, widths=[1.3, 1.6, 0.5, 3.2], fs=9)

    # ---------------- 4
    h("4. ER Model", 1)
    para("Figure 1 shows the whole design as one ER diagram in Chen notation: "
         "rectangles are entities, diamonds are relationships and ellipses are "
         "attributes. An underlined attribute is the key. The 1, M and N marks on "
         "each line say how many rows can take part, and a double line means taking "
         "part is compulsory, so for example every order must belong to a customer.")
    figure("er_model.png", "Figure 1. ER model with 9 entities and 12 "
                           "relationships")
    para("Two things in Figure 1 are worth pointing out before the EER model "
         "changes them. MANAGES is the only 1:1 relationship in the design: a "
         "warehouse has exactly one manager, and an employee manages at most one "
         "warehouse, so both sides are marked 1 and the warehouse side is drawn "
         "with a double line because every warehouse must have one. And INSPECTION "
         "has to be joined twice, once to SUPPLIER and once to PRODUCT, because "
         "plain ER has no way of saying that an inspection is about a supplier and "
         "a product taken together. Section 5.2 replaces those two links with one.")

    # ---------------- 5
    h("5. EER Model", 1)
    para("An EER model is the same ER model with a few extra ideas added, so it can "
         "describe situations plain ER cannot. Figure 2 is the identical design from "
         "Figure 1, with everything EER contributes drawn shaded. Two of those "
         "ideas need a word on direction. Specialization works downwards: EMPLOYEE "
         "already existed and is split into MANAGER and DRIVER. Generalization "
         "works upwards: SHIPMENT and GOODS_RETURN already existed as separate "
         "entities, and noticing what they share creates GOODS_MOVEMENT above "
         "them. The two produce the same shape on paper and map to tables the same "
         "way, but they are arrived at from opposite ends.")
    figure("eer_model.png", "Figure 2. The same model, extended with EER features")
    para("Figure 2 uses several pieces of notation that Figure 1 does not need. A "
         "double ellipse is a multivalued attribute, so SUPPLIER.Phone is drawn that "
         "way. A composite attribute is drawn as an ellipse with its parts hanging "
         "below it, which is how CUSTOMER.Address carries Street, City and Pincode. "
         "A dashed ellipse is a derived attribute, used for ORDERS.TotalAmt. A "
         "rectangle drawn around a whole relationship marks an aggregation, which "
         "section 5.2 covers. A weak entity has a double border and reaches its owner "
         "through a double diamond, which is why ORDER_ITEM and HAS are drawn that "
         "way. The generalization at the foot of the figure uses the same circle "
         "and double line as a specialization, because on paper the two constructs "
         "cannot be told apart.")
    para("Each of the three hierarchies carries three marks, and every mark answers "
         "a different question:")
    bullets([
        "The letter in the circle says whether the subclasses may overlap. A d means "
        "disjoint, so a row joins at most one subclass: an employee is a manager or "
        "a driver and never both, and a movement is a shipment or a return. An o "
        "means overlapping, so a row may join more than one: a product can be both "
        "perishable and hazardous, and its ProductID then appears in both subclass "
        "tables.",
        "The line running into the circle says whether every row has to join a "
        "subclass at all. A double line means total, which is why EMPLOYEE and "
        "GOODS_MOVEMENT have one. A single line means partial, which is why PRODUCT "
        "has one: most products are neither perishable nor hazardous.",
        "The subset symbol on each line from a subclass up to the circle points at "
        "the superclass and says that every row of the subclass is also a row of the "
        "superclass. It is the reason the subclass tables share the superclass key "
        "rather than inventing one of their own.",
    ])
    para("Because SHIPMENT and GOODS_RETURN now sit under one superclass, the two "
         "links each of them had to ORDERS and to WAREHOUSE become one link each, "
         "MOVEMENT_FOR and HANDLED_AT, and their two keys collapse into the single "
         "key MovementID. What stays with a subclass is what only that subclass "
         "has: DeliveryDate on SHIPMENT, Reason and RefundAmt on GOODS_RETURN.")
    h("5.1 What EER Added and Why", 2)
    table(["EER feature", "Where", "Why it is needed"],
          EER_ADDED, widths=[1.4, 1.5, 3.7], fs=9)
    h("5.2 Aggregation", 2)
    para("Aggregation is drawn on its own in Figure 3, because the box it needs "
         "would sit on top of half of Figure 2. The box encloses SUPPLIER, "
         "SUPPLIES and PRODUCT and says: treat that whole arrangement as one "
         "object. INSPECTION then joins to the box rather than to either entity "
         "inside it, which is exactly the statement plain ER could not make. When "
         "this maps to tables the aggregate is the SUPPLIES table, and INSPECTION "
         "carries SupplierID and ProductID together as one foreign key into it, so "
         "an inspection can only ever name a pair that really is supplied.")
    figure("eer_aggregation.png", "Figure 3. The aggregation: INSPECTION joins the "
                                  "boxed SUPPLIES arrangement, not its members",
           width=5.0)
    h("5.3 The Weak Entity", 2)
    para("The weak entity is worth singling out. ORDER_ITEM is drawn with a double "
         "border because it cannot exist on its own, and the diamond joining it to "
         "ORDERS is drawn as a double diamond to show that this is the relationship "
         "that gives it its identity. Its partial key, ItemNo, is underlined with a "
         "dashed line because it is only unique within one order.", italic=True)

    # ---------------- 6
    h("6. Converting the EER Model into Tables", 1)
    para("The EER model is converted into relational tables using the standard "
         "mapping rules:")
    bullets(MAPPING_RULES)
    para("Applying these rules gives the tables that the rest of the report works "
         "on. Two of them leave a mark worth noting. MANAGES puts ManagerID on "
         "WAREHOUSE with a UNIQUE constraint, which together with EMPLOYEE holding "
         "WarehouseID makes the two tables reference each other, so one of the two "
         "columns has to be filled in after both rows exist. And the aggregation "
         "gives INSPECTION a foreign key of two columns at once, SupplierID with "
         "ProductID, pointing at SUPPLIES rather than at SUPPLIER and PRODUCT "
         "separately. The finished set of tables is listed in section 10.")

    # ---------------- 7
    h("7. Sample Data and Functional Dependencies", 1)
    para("A functional dependency, written X -> Y, means that if you know X then Y "
         "is fixed. For example CustID -> CName: once the customer's ID is known, "
         "the name is decided. Every normal form is a rule about which dependencies "
         "a table is allowed to contain, so the dependencies have to be settled "
         "before any splitting starts.")
    para("Sample rows are what make that possible. Data can never prove a "
         "dependency, because a later row may still break it, but a single pair of "
         "rows can disprove one: if two rows agree on X and disagree on Y, then "
         "X -> Y is dead. So each relation below is filled with at least five rows, "
         "every candidate dependency is checked against them, and the survivors are "
         "kept only if they also make sense as a rule of the business.")
    figure("sample_data.png", "Figure 4. Sample rows used for the dependency "
                              "analysis. The seven ORDER_LINE rows are numbered 1 to "
                              "7 from the top.")
    para("ORDER_LINE holds every order and product fact in one relation, which is "
         "the state the design is in before normalization. DELIVERY_DUTY records "
         "which driver took an order out in which city, under two rules the "
         "transport team works to: each driver covers exactly one city, and within a "
         "city an order always goes to the same driver.")

    h("7.1 Dependencies That Hold", 2)
    table(["Ref", "Dependency", "Evidence in the sample rows"],
          [(a, b, c) for a, b, c, _ in FD_HOLD], widths=[0.5, 2.3, 3.8], fs=9)

    h("7.2 Dependencies That Do Not Hold", 2)
    para("Each of these is ruled out by two rows that agree on the left-hand side "
         "and disagree on the right. They are listed because they are the ones a "
         "reader might otherwise assume, and because ruling them out is what makes "
         "the candidate key in 7.3 the only one.")
    table(["Ref", "Rejected dependency", "Counterexample"],
          FD_REJECT, widths=[0.5, 2.3, 3.8], fs=9)

    h("7.3 Candidate Keys", 2)
    para("A candidate key is a set of columns whose closure is the whole relation. "
         "For ORDER_LINE, start with {OrderID, ProductID}:")
    bullets([
        "F1 adds OrderDate, Status and CustID.",
        "F2 then adds CName, Street and City.",
        "F3 then adds State.",
        "F4 adds PName, Category and Price.",
        "F5 adds Qty.",
    ])
    para("The closure is now all thirteen columns, so {OrderID, ProductID} is a "
         "key. Neither half works alone: OrderID cannot fix ProductID (N2) and "
         "ProductID cannot fix Qty (N1), so nothing smaller is a key and this is the "
         "only candidate key of ORDER_LINE.")
    para("DELIVERY_DUTY has two candidate keys. {OrderID, City} determines DriverID "
         "by F7, and {OrderID, DriverID} determines City by F6, so both close over "
         "the whole relation. That second key is what lets DELIVERY_DUTY pass 3NF "
         "while still failing BCNF in section 8.5. SUPPLIER_PHONE also has two: the "
         "pair {SupplierID, Phone}, and Phone on its own by F8.")

    # ---------------- 8
    h("8. Normalization", 1)
    para("Normalization starts from one badly designed table and splits it, one rule "
         "at a time, until no table can store the same fact twice. Each step below "
         "names the dependency from section 7 that forces the split, and the figure "
         "beside it shows the tables that result, with the sample rows carried "
         "through so the effect of the split can be read off directly.")
    for i, (name, heading, text, verdict, _) in enumerate(NORM_STEPS):
        h(heading, 2)
        para(text)
        para(verdict, italic=True, size=10.5)
        figure(name, "Figure " + str(5 + i) + ". " + heading.split("  ", 1)[1])

    # ---------------- 9
    h("9. The Decomposition, End to End", 1)
    para("Figure 10 puts the five stages side by side and follows every relation "
         "through them, with all of its columns and its key at each stage. Reading "
         "left to right, one relation becomes four at 1NF, ORDER_ITEM sheds a "
         "PRODUCT table at 2NF, ORDER_MASTER becomes three tables at 3NF, and "
         "DELIVERY_DUTY becomes two at BCNF. A relation drawn in a bold box was "
         "created or changed at that stage, and a thin box is carried forward "
         "untouched.")
    figure("decomp_tree.png", "Figure 10. Decomposition from UNF to BCNF, showing "
                              "every column at every stage")
    para("Eight relations come out of this decomposition. The remaining thirteen "
         "tables of the final design (SUPPLIER, SUPPLIES, INSPECTION, "
         "PERISHABLE_PRODUCT, HAZARDOUS_PRODUCT, WAREHOUSE, EMPLOYEE, MANAGER, "
         "DRIVER, STOCK, GOODS_MOVEMENT, SHIPMENT and GOODS_RETURN) come straight "
         "from the EER mapping in section 6, and each of them was already in BCNF. "
         "Eight plus thirteen is the twenty-one tables of section 10.")

    # ---------------- 10
    h("10. Final Schema", 1)
    para("The finished design has 21 tables and every one of them is in BCNF, which "
         "also means every one is in 1NF, 2NF and 3NF. In Figure 11 the tables are "
         "grouped under three headings, master data and supply, warehouse and "
         "staff, and orders and movements. A column name is underlined where it "
         "forms part of the primary key and italicised where it is a foreign key, "
         "so a column that is both is underlined and italic.")
    figure("final_schema.png", "Figure 11. Final relational schema, 21 tables in "
                               "BCNF")
    h("10.1 Keys of Every Table", 2)
    table(["Table", "Primary key", "Foreign keys"], SCHEMA,
          widths=[2.1, 2.2, 2.6], fs=9.5)
    h("10.2 Why the Design Is Now Safe", 2)
    bullets([
        "Every fact is stored in exactly one place, so two rows can never disagree.",
        "Correcting a product name, a customer address or a city's state is a "
        "single-row edit.",
        "A supplier, product or customer can be added before it appears on any "
        "order.",
        "Deleting an order cannot wipe out the customer who placed it.",
        "TotalAmt is calculated from the order lines rather than stored, so it can "
        "never fall out of step with them.",
        "A movement is recorded once, in GOODS_MOVEMENT, whether it is a shipment "
        "or a return, so neither kind can go missing from a count of what left the "
        "warehouse.",
        "An inspection cannot name a supplier and a product that are not actually "
        "paired, because its foreign key is the pair.",
        "A product that is both perishable and hazardous is one row in PRODUCT with "
        "one row in each subclass table, so neither set of handling rules has to be "
        "copied into the other.",
    ])
    para("Two tables are worth a second look. WAREHOUSE has two candidate keys, "
         "WarehouseID and ManagerID, because MANAGES is 1:1 and ManagerID is "
         "therefore UNIQUE. Both determinants are candidate keys, so the table is "
         "still in BCNF. GOODS_MOVEMENT, SHIPMENT and GOODS_RETURN each have "
         "MovementID as their only key and no other determinant, which is what "
         "generalization gives for free: the shared columns sit in one table, so "
         "there is nowhere for a second copy of them to disagree.")
    para("One trade-off is worth recording. Splitting DELIVERY_DUTY into "
         "DRIVER_CITY and ORDER_DRIVER means F7, {OrderID, City} -> DriverID, can no "
         "longer be checked inside a single table. It will be enforced by a trigger "
         "in DA2. This is the normal price of BCNF, and it is worth paying here "
         "because the alternative allows a driver to be recorded against two "
         "different cities at once.")

    # ---------------- 11
    h("11. Conclusion and Next Steps", 1)
    para("DA1 delivered an ER model of nine entities and twelve relationships, the "
         "same model extended into an EER model with specialization, both total and "
         "disjoint and partial and overlapping, generalization, aggregation, a weak "
         "entity, a 1:1 relationship made precise, and "
         "multivalued, composite and derived attributes, a dependency analysis "
         "carried out on sample rows, and a full normalization from UNF to BCNF in "
         "which every step is shown with the tables before and after. The result is "
         "21 tables, all in BCNF.")
    para("DA2, due 18 September 2026: ", bold=True)
    d.paragraphs[-1].add_run(
        "create the 21 tables in SQL with their keys and constraints, load sample "
        "data, add the trigger that enforces the dependency lost at BCNF, and write "
        "PL/SQL procedures for receiving stock and dispatching an order.")
    para("DA3, due 23 October 2026: ", bold=True)
    d.paragraphs[-1].add_run(
        "rebuild the same project on one modern database technology. The current "
        "plan is a graph database (Neo4j), because tracing a delivery back through "
        "shipment, order, warehouse and supplier is a chain of hops that a graph "
        "query expresses in one line, whereas SQL needs a fresh join for every "
        "extra step.")

    out = os.path.join(ROOT, "docs", "DA1_Report.docx")
    d.save(out)
    return stamp_office_metadata(out, "Microsoft Office Word")


# =====================================================================  PPTX
def build_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    BLANK = prs.slide_layouts[6]
    W = prs.slide_width

    cp = prs.core_properties
    cp.title = TITLE
    cp.subject = "Database Management Systems, Digital Assignment 1"
    cp.author = DOC_AUTHOR
    cp.last_modified_by = DOC_AUTHOR
    cp.category = ""
    cp.comments = ""
    cp.keywords = ""
    cp.created = CREATED
    cp.modified = MODIFIED
    cp.revision = 9

    def textbox(slide, l, t, w, h, text, size, color=PBLACK, bold=False,
                align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(PInches(l), PInches(t), PInches(w), PInches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = PPt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        f.name = "Calibri"
        return tb

    def rule(slide, y, l=0.5, w=12.4):
        ln = slide.shapes.add_shape(1, PInches(l), PInches(y), PInches(w),
                                    PInches(0.02))
        ln.fill.solid()
        ln.fill.fore_color.rgb = PBLACK
        ln.line.fill.background()
        ln.shadow.inherit = False
        return ln

    def header(slide, title, subtitle=None, top=0.32):
        textbox(slide, 0.5, top, 12.4, 0.62, title, 26, PBLACK, bold=True)
        rule(slide, top + 0.72)
        if subtitle:
            textbox(slide, 0.5, top + 0.80, 12.4, 0.42, subtitle, 13.5, PGREY,
                    italic=True)
        return top + (1.34 if subtitle else 0.92)

    def bullets_slide(title, items, subtitle=None, size=17):
        s = prs.slides.add_slide(BLANK)
        top = header(s, title, subtitle)
        tb = s.shapes.add_textbox(PInches(0.7), PInches(top), PInches(12.0),
                                  PInches(7.2 - top))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            sub = it.startswith("  ")
            r = p.add_run()
            r.text = ("-  " if sub else "•  ") + it.strip()
            r.font.size = PPt(size - 2 if sub else size)
            r.font.color.rgb = PGREY if sub else PBLACK
            r.font.name = "Calibri"
            p.space_after = PPt(9)
            if sub:
                p.level = 1
        return s

    def image_slide(title, image, subtitle=None):
        s = prs.slides.add_slide(BLANK)
        top = header(s, title, subtitle, top=0.26)
        iw, ih = Image.open(image).size
        maxw, maxh = 12.6, 7.30 - top
        scale = min(maxw / iw, maxh / ih)
        dw, dh = iw * scale, ih * scale
        s.shapes.add_picture(image, PInches((13.333 - dw) / 2),
                             PInches(top + (maxh - dh) / 2), PInches(dw), PInches(dh))
        return s

    def table_slide(title, headers, rows, widths, subtitle=None, fs=12):
        s = prs.slides.add_slide(BLANK)
        top = header(s, title, subtitle)
        rowh = 0.44
        est = rowh * (len(rows) + 1)
        top += max(0.0, (7.25 - top - est) / 2)
        total = sum(widths)
        shape = s.shapes.add_table(len(rows) + 1, len(headers),
                                   PInches((13.333 - total) / 2), PInches(top),
                                   PInches(total), PInches(est))
        tbl = shape.table
        for i, w in enumerate(widths):
            tbl.columns[i].width = PInches(w)
        for i, hd in enumerate(headers):
            c = tbl.cell(0, i)
            c.text = hd
            pr = c.text_frame.paragraphs[0]
            pr.runs[0].font.size = PPt(fs)
            pr.runs[0].font.bold = True
            pr.runs[0].font.color.rgb = PBLACK
            c.fill.solid()
            c.fill.fore_color.rgb = PHEAD
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                c = tbl.cell(ri, ci)
                c.text = str(val)
                pr = c.text_frame.paragraphs[0]
                pr.runs[0].font.size = PPt(fs - 1)
                pr.runs[0].font.color.rgb = PBLACK
                c.fill.solid()
                c.fill.fore_color.rgb = PZEBRA if ri % 2 else PWHITE
                c.vertical_anchor = MSO_ANCHOR.MIDDLE
        return s

    # ---------------- title
    s = prs.slides.add_slide(BLANK)
    textbox(s, 0.8, 2.05, 11.7, 1.7, TITLE, 40, PBLACK, bold=True,
            align=PP_ALIGN.CENTER)
    rule(s, 3.62, l=3.4, w=6.5)
    textbox(s, 0.8, 3.80, 11.7, 0.6, "Digital Assignment 1, Database Management "
            "Systems", 20, PGREY, align=PP_ALIGN.CENTER)
    textbox(s, 0.8, 4.45, 11.7, 0.6, "ER / EER Model and Normalization up to BCNF",
            19, PBLACK, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 0.8, 5.75, 11.7, 0.9,
            "   |   ".join(f"{n} ({r})" for n, r in TEAM), 15, PBLACK, bold=True,
            align=PP_ALIGN.CENTER)

    # ---------------- agenda
    bullets_slide("Agenda", [
        "The problem, objectives and scope",
        "Entities, attributes and relationships",
        "The ER model",
        "The EER model: specialization, generalization and aggregation",
        "Turning the model into tables",
        "Sample data and the functional dependencies it reveals",
        "Normalization: UNF to 1NF to 2NF to 3NF to BCNF",
        "The decomposition end to end, the final schema, and what comes next",
    ])

    # ---------------- problem
    bullets_slide("1. The Problem", [
        "A warehouse must always know: what stock is where, who supplies it, and "
        "where each order has reached.",
        "In practice those three answers live in three separate places, and nothing "
        "keeps them in agreement.",
        "So the same fact gets stored twice: a customer's address on every order, a "
        "product's name on every order line.",
        "Correct one copy, miss another, and the data contradicts itself.",
    ], subtitle="the trouble is not missing data, it is the same data stored more "
                "than once")

    bullets_slide("1. What That Duplication Costs", [
        "Insert problem: a new customer cannot be recorded until they place an "
        "order.",
        "Delete problem: deleting a customer's last order erases the customer.",
        "Update problem: renaming a product means editing many rows, and missing one "
        "leaves two answers.",
        "Unreliable stock: the quantity on hand is kept in more than one place.",
        "These are the insert, delete and update anomalies. Normalization removes "
        "them.",
    ])

    bullets_slide("1. Objectives and Scope", [
        "Identify the entities, attributes and relationships.",
        "Draw the ER model, then extend the same model into an EER model.",
        "Convert the model into relational tables.",
        "Work out the functional dependencies from sample rows.",
        "Normalize every table from UNF up to BCNF, showing each step.",
        "In scope: 9 entities, namely warehouse, product, supplier, customer, "
        "orders, shipment, return, inspection and employee.",
        "Out of scope: SQL and PL/SQL (DA2), and a modern database (DA3).",
    ])

    # ---------------- entities & relationships
    table_slide("2. The Nine Entities",
                ["Entity", "Key attribute", "Other attributes"],
                ENTITIES, [2.4, 2.5, 6.6],
                subtitle="the key attribute identifies one row, and is underlined "
                         "in every diagram", fs=12)

    table_slide("3. The Twelve Relationships",
                ["Relationship", "Between", "Type"],
                [(r[0], r[1], r[2]) for r in RELATIONSHIPS], [3.4, 4.6, 1.4],
                subtitle="1:1 means one row on each side, 1:N one to many, M:N many "
                         "on both sides", fs=11.5)

    # ---------------- ER / EER
    image_slide("4. The ER Model", img("er_model.png"),
                subtitle="MANAGES is the one 1:1 relationship; INSPECTION needs two "
                         "links because plain ER cannot point at a pair")
    image_slide("5. The EER Model, the Same Design Extended", img("eer_model.png"))

    bullets_slide("5. Specialization and Generalization", [
        "Specialization works downwards, from a type we already had to its kinds.",
        "  EMPLOYEE existed; splitting it gives MANAGER (Level, Bonus) and DRIVER "
        "(LicenseNo, Expiry). Total and disjoint.",
        "  PRODUCT existed; PERISHABLE_PRODUCT and HAZARDOUS_PRODUCT are partial "
        "(most products are neither) and overlapping (a vaccine is both), so that "
        "circle is marked o and not d.",
        "Generalization works upwards, from types we already had to what they share.",
        "  SHIPMENT and GOODS_RETURN were both in the ER model. Both are one "
        "order's goods moving on a date with a status, so GOODS_MOVEMENT is created "
        "above them.",
        "  Their two links to ORDERS and WAREHOUSE become one each, and ShipmentID "
        "and ReturnID collapse into MovementID.",
        "Same shape on paper, same mapping to tables, opposite directions of "
        "reasoning.",
    ], size=15, subtitle="the difference is the direction, not the drawing")

    bullets_slide("5. Aggregation", [
        "An inspection is of one supplier's supply of one product.",
        "Not of the supplier, and not of the product, which is why the ER model "
        "needed two separate links and still could not say it.",
        "So the whole SUPPLIER, SUPPLIES, PRODUCT arrangement is boxed and treated "
        "as one object, and INSPECTION joins the box.",
        "In tables: the aggregate is the SUPPLIES table, and INSPECTION carries "
        "SupplierID with ProductID as a single foreign key into it.",
        "An inspection therefore cannot name a pair that is not actually supplied.",
    ], subtitle="a relationship treated as one higher-level object")

    image_slide("5. The Aggregation on Its Own", img("eer_aggregation.png"))

    table_slide("5. What EER Added, and Why",
                ["EER feature", "Where"],
                [(a.replace("\n", " "), b.replace("\n", " "))
                 for a, b, _ in EER_ADDED], [4.6, 5.6],
                subtitle="each addition describes something plain ER cannot",
                fs=13)

    bullets_slide("5. The Weak Entity", [
        "ORDER_ITEM cannot exist on its own, an order line means nothing without "
        "its order.",
        "ItemNo (1, 2, 3 and so on) is only unique inside one order, so it is a "
        "partial key, shown with a dashed underline.",
        "Its real key is therefore OrderID + ItemNo.",
        "Drawn with a double border, and the diamond joining it to ORDERS is a "
        "double diamond, meaning that relationship gives it its identity.",
        "It replaces the M:N CONTAINS relationship from the ER model.",
    ], subtitle="why ORDER_ITEM is drawn differently from every other box")

    # ---------------- mapping
    bullets_slide("6. Turning the Model into Tables", MAPPING_RULES,
                  subtitle="the standard mapping rules", size=16)

    # ---------------- dependencies
    bullets_slide("7. Why Sample Data Comes First", [
        "A functional dependency X -> Y means: if you know X, then Y is fixed.",
        "  Example, CustID -> CName: once the customer ID is known, the name is "
        "decided.",
        "Each normal form is a rule about which dependencies a table may contain, so "
        "the dependencies have to be settled before any splitting starts.",
        "Data cannot prove a dependency, but two rows can kill one: agree on X, "
        "disagree on Y, and X -> Y is gone.",
        "So every relation gets at least five rows, and every candidate dependency "
        "is checked against them.",
    ], subtitle="the dependencies drive every split that follows")

    image_slide("7. The Sample Rows", img("sample_data.png"),
                subtitle="ORDER_LINE rows are numbered 1 to 7 from the top")

    table_slide("7. Dependencies That Hold", ["Ref", "Dependency", "Read off"],
                [(a, b, d) for a, b, _, d in FD_HOLD], [1.0, 4.0, 7.2], fs=11.5,
                subtitle="each one survives every sample row")

    table_slide("7. Dependencies the Data Rules Out",
                ["Ref", "Rejected", "Counterexample"],
                [(a, b, c) for a, b, c in FD_REJECT], [1.0, 3.4, 7.8], fs=11,
                subtitle="two rows that agree on the left and disagree on the right")

    bullets_slide("7. Candidate Keys", [
        "ORDER_LINE: start from {OrderID, ProductID} and close it.",
        "  F1 adds OrderDate, Status, CustID. F2 adds CName, Street, City. F3 adds "
        "State. F4 adds PName, Category, Price. F5 adds Qty.",
        "  All thirteen columns are reached, and neither half works alone (N1, N2), "
        "so it is the only candidate key.",
        "DELIVERY_DUTY: two candidate keys, {OrderID, City} by F7 and "
        "{OrderID, DriverID} by F6.",
        "  That second key is exactly why it passes 3NF and still fails BCNF.",
        "SUPPLIER_PHONE: {SupplierID, Phone}, and Phone alone by F8.",
    ], size=16)

    # ---------------- normalization
    bullets_slide("8. The Four Rules", [
        "1NF: every cell holds one value, and repeating groups move out.",
        "2NF: no column depends on only part of the key.",
        "3NF: no ordinary column decides another ordinary column.",
        "BCNF: the left side of every dependency must be a key.",
        "Each step below names the dependency from section 7 that forces the split.",
    ], subtitle="start with one bad table, then split it one rule at a time")

    for name, heading, _, _, points in NORM_STEPS:
        _, rest = heading.split("  ", 1)
        short = rest.split("(")[-1].rstrip(")") if "(" in rest else rest
        bullets_slide("8. " + rest, points, size=16)
        image_slide("8. The Tables after " + short, img(name))

    # ---------------- tree and final
    image_slide("9. The Decomposition, End to End", img("decomp_tree.png"),
                subtitle="every relation, every column, at every stage")

    image_slide("10. Final Schema, 21 Tables, All in BCNF", img("final_schema.png"))

    bullets_slide("10. Why the Design Is Now Safe", [
        "Every fact is stored in exactly one place, so two rows can never disagree.",
        "Correcting a product name, an address or a city's state is a single-row "
        "edit.",
        "A supplier, product or customer can exist before appearing on any order.",
        "Deleting an order cannot wipe out the customer who placed it.",
        "TotalAmt is calculated, not stored, so it can never fall out of step.",
        "A movement is recorded once whether it is a shipment or a return.",
        "WAREHOUSE has two candidate keys now, WarehouseID and ManagerID, because "
        "MANAGES is 1:1. Both are keys, so it is still in BCNF.",
        "One trade-off: F7, {OrderID, City} -> DriverID, now needs a trigger, which "
        "DA2 will add.",
    ], size=16)

    bullets_slide("11. What Comes Next", [
        "DA2, due 18 September 2026",
        "  Create the 21 tables in SQL with keys and constraints, load sample data, "
        "add the BCNF trigger, and write PL/SQL for receiving stock and dispatching "
        "orders.",
        "DA3, due 23 October 2026",
        "  Rebuild on a graph database (Neo4j). Tracing a delivery back through "
        "shipment, order, warehouse and supplier is a chain of hops, one line in a "
        "graph query but a fresh join for every step in SQL.",
    ], size=17)

    # ---------------- closing
    s = prs.slides.add_slide(BLANK)
    textbox(s, 0.8, 2.9, 11.7, 1.2, "Thank You", 44, PBLACK, bold=True,
            align=PP_ALIGN.CENTER)
    rule(s, 4.15, l=4.4, w=4.5)
    textbox(s, 0.8, 4.35, 11.7, 0.7, "9 entities, 12 relationships, 21 tables, all "
            "in BCNF", 19, PGREY, italic=True, align=PP_ALIGN.CENTER)
    textbox(s, 0.8, 5.7, 11.7, 0.7,
            "   |   ".join(f"{n} ({r})" for n, r in TEAM), 14, PBLACK,
            align=PP_ALIGN.CENTER)

    out = os.path.join(ROOT, "presentation", "DA1_Presentation.pptx")
    n = len(prs.slides._sldIdLst)
    prs.save(out)
    stamp_office_metadata(out, "Microsoft Office PowerPoint",
                          fmt="Widescreen")
    return out, n


if __name__ == "__main__":
    for f in ["er_model", "eer_model", "eer_aggregation", "sample_data",
              "norm_unf", "norm_1nf", "norm_2nf", "norm_3nf", "norm_bcnf",
              "decomp_tree", "final_schema"]:
        assert os.path.exists(img(f + ".png")), \
            f"missing {f}.png, run diagrams/make_diagrams.py"
    print("DOCX:", build_docx())
    path, n = build_pptx()
    print("PPTX:", path, f"({n} slides)")
